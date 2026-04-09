"""
Generate PowerPoint presentation for midterm project.
Uses python-pptx to create a professional slide deck.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

REPORT_DIR = os.path.dirname(__file__)
FIG_DIR = os.path.join(REPORT_DIR, "figures")
COW_IMG = r"C:\Users\Samih\Pictures\Screenshots\Screenshot 2026-04-08 221810.png"

# Use PNG versions of figures (generated alongside PDFs by generate_figures.py)

# Brand colors
DARK_BG = RGBColor(0x1B, 0x1B, 0x2F)       # deep navy
ACCENT = RGBColor(0x00, 0x96, 0xC7)          # teal blue
ACCENT2 = RGBColor(0xE0, 0x6C, 0x00)         # warm orange
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF5)
MID_GRAY = RGBColor(0x90, 0x90, 0xA0)
TEXT_DARK = RGBColor(0x22, 0x22, 0x33)
HIGHLIGHT = RGBColor(0xFF, 0xD6, 0x00)        # gold
RED_ACCENT = RGBColor(0xE0, 0x40, 0x40)
GREEN_ACCENT = RGBColor(0x40, 0xB0, 0x60)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


def add_bg(slide, color=DARK_BG):
    """Fill slide background with solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, color, alpha=None):
    """Add a rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    """Add a text box with styled text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_slide_content(slide, left, top, width, height, bullets,
                              font_size=16, color=WHITE, spacing=Pt(6)):
    """Add bulleted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, bold_part) in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = spacing
        if bold_part:
            run_b = p.add_run()
            run_b.text = bold_part
            run_b.font.size = Pt(font_size)
            run_b.font.color.rgb = ACCENT
            run_b.font.bold = True
            run_b.font.name = "Calibri"
            run_r = p.add_run()
            run_r.text = text
            run_r.font.size = Pt(font_size)
            run_r.font.color.rgb = color
            run_r.font.name = "Calibri"
        else:
            run = p.add_run()
            run.text = text
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.name = "Calibri"
    return txBox


def section_header(slide, number, title):
    """Add a colored accent bar + section number + title at top."""
    add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), ACCENT)
    add_text_box(slide, Inches(0.7), Inches(0.35), Inches(1.5), Inches(0.5),
                 f"{number:02d}", font_size=14, color=MID_GRAY, bold=False)
    add_text_box(slide, Inches(0.7), Inches(0.65), Inches(11), Inches(0.8),
                 title, font_size=32, color=WHITE, bold=True)


def add_image_safe(slide, path, left, top, width=None, height=None):
    """Add image if it exists."""
    if os.path.exists(path):
        kwargs = {}
        if width:
            kwargs["width"] = width
        if height:
            kwargs["height"] = height
        slide.shapes.add_picture(path, left, top, **kwargs)
        return True
    else:
        add_text_box(slide, left, top, Inches(4), Inches(1),
                     f"[Image not found: {os.path.basename(path)}]",
                     font_size=12, color=RED_ACCENT)
        return False


# =============================================================================
# SLIDE 1: TITLE
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.12), ACCENT)
add_shape(slide, Inches(0), H - Inches(0.12), W, Inches(0.12), ACCENT)

add_text_box(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(1.2),
             "Deep Learning for Vessel Segmentation\nin Cerebral CTA",
             font_size=40, color=WHITE, bold=True)

add_text_box(slide, Inches(1.5), Inches(3.2), Inches(10), Inches(0.8),
             "A Controlled Ablation of Topology-Aware Loss Functions",
             font_size=24, color=ACCENT)

add_shape(slide, Inches(1.5), Inches(4.3), Inches(3), Inches(0.04), ACCENT)

add_text_box(slide, Inches(1.5), Inches(4.7), Inches(6), Inches(0.5),
             "Samih Tharwat Amer", font_size=20, color=WHITE, bold=True)
add_text_box(slide, Inches(1.5), Inches(5.2), Inches(6), Inches(0.5),
             "Whiting School of Engineering  |  Johns Hopkins University",
             font_size=16, color=MID_GRAY)
add_text_box(slide, Inches(1.5), Inches(5.7), Inches(6), Inches(0.5),
             "Midterm Project  |  April 2026", font_size=14, color=MID_GRAY)


# =============================================================================
# SLIDE 2: CLINICAL MOTIVATION
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_header(slide, 1, "Clinical Motivation")

add_bullet_slide_content(slide, Inches(0.7), Inches(1.7), Inches(5.5), Inches(4.5), [
    (" are a leading cause of death and disability worldwide", "Cerebrovascular diseases"),
    (" from CTA is essential for diagnosis, surgical navigation, and treatment planning", "Vessel segmentation"),
    (" — the arterial ring at the brain's base — provides collateral blood flow", "Circle of Willis (CoW)"),
    (" (Acom, Pcom) are thin but clinically critical for aneurysm assessment", "Communicating arteries"),
], font_size=17)

add_image_safe(slide, COW_IMG, Inches(6.8), Inches(1.5), width=Inches(5.8))


# =============================================================================
# SLIDE 3: THE PROBLEM
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_header(slide, 2, "The Problem: Topological Failures")

add_bullet_slide_content(slide, Inches(0.7), Inches(1.7), Inches(7), Inches(2.5), [
    (" treats all voxels equally — no connectivity incentive", "Dice + Cross-Entropy loss"),
    (" can have high Dice (0.87) but broken topology — severed vessels, fragments", "A segmentation"),
    (" hide failures on thin communicating arteries", "Aggregate metrics"),
], font_size=18)

# Big callout box
box = add_shape(slide, Inches(1.5), Inches(4.0), Inches(10), Inches(2.2), RGBColor(0x25, 0x25, 0x40))
add_text_box(slide, Inches(2), Inches(4.2), Inches(9), Inches(0.6),
             "\"The model missed a significant number of vessels in clinical use,",
             font_size=18, color=ACCENT, bold=False)
add_text_box(slide, Inches(2), Inches(4.7), Inches(9), Inches(0.6),
             "despite the validation DSC looking reasonable on paper.\"",
             font_size=18, color=ACCENT, bold=False)

add_text_box(slide, Inches(2), Inches(5.4), Inches(9), Inches(0.5),
             "Global DSC 0.87  —  but R-Pcom DSC 0.42, L-ACA DSC 0.46",
             font_size=16, color=RED_ACCENT, bold=True)


# =============================================================================
# SLIDE 4: HYPOTHESIS
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_header(slide, 3, "Hypothesis & Approach")

# Hypothesis callout
box = add_shape(slide, Inches(0.7), Inches(1.8), Inches(11.5), Inches(1.5), RGBColor(0x00, 0x3D, 0x5B))
add_text_box(slide, Inches(1.0), Inches(1.85), Inches(11), Inches(0.4),
             "HYPOTHESIS", font_size=12, color=ACCENT, bold=True)
add_text_box(slide, Inches(1.0), Inches(2.2), Inches(11), Inches(1.0),
             "Topology-aware loss functions will disproportionately improve\nsegmentation of thin communicating arteries (Acom, Pcom)",
             font_size=20, color=WHITE, bold=True)

# Three loss configs
configs = [
    ("1", "Dice + CE", "Baseline — overlap only", ACCENT),
    ("2", "Dice + CE + clDice", "Centerline matching (CVPR 2021)", ACCENT2),
    ("3", "Dice + CE + Skeleton Recall", "Centerline recall (ECCV 2024)", GREEN_ACCENT),
]
for i, (num, name, desc, color) in enumerate(configs):
    left = Inches(0.7 + i * 4.1)
    box = add_shape(slide, left, Inches(3.8), Inches(3.8), Inches(2.0), RGBColor(0x25, 0x25, 0x40))
    add_text_box(slide, left + Inches(0.2), Inches(3.9), Inches(3.4), Inches(0.5),
                 f"Config {num}", font_size=12, color=MID_GRAY, bold=False)
    add_text_box(slide, left + Inches(0.2), Inches(4.25), Inches(3.4), Inches(0.6),
                 name, font_size=18, color=color, bold=True)
    add_text_box(slide, left + Inches(0.2), Inches(4.85), Inches(3.4), Inches(0.6),
                 desc, font_size=14, color=MID_GRAY)

add_text_box(slide, Inches(0.7), Inches(6.2), Inches(11), Inches(0.5),
             "Controlled ablation: same architecture, same data, same hardware — only the loss changes",
             font_size=15, color=MID_GRAY, bold=False)


# =============================================================================
# SLIDE 5: DATASETS
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_header(slide, 4, "Datasets")

# TopCoW box
box = add_shape(slide, Inches(0.7), Inches(1.7), Inches(5.8), Inches(4.8), RGBColor(0x25, 0x25, 0x40))
add_text_box(slide, Inches(1.0), Inches(1.8), Inches(5), Inches(0.5),
             "TopCoW 2024", font_size=22, color=ACCENT, bold=True)
add_text_box(slide, Inches(1.0), Inches(2.3), Inches(5), Inches(0.4),
             "MICCAI Challenge Dataset", font_size=13, color=MID_GRAY)
add_bullet_slide_content(slide, Inches(1.0), Inches(2.8), Inches(5.2), Inches(3.5), [
    ("", "125 CTA volumes"),
    (" Circle of Willis vessel classes", "13"),
    (": BA, ICA, MCA, PCA, ACA, Acom, Pcom, 3rd-A2", "Classes"),
    (" train / 25 validation (80/20 split)", "100"),
    (" for from-scratch training", "Used"),
], font_size=15)

# TopBrain box
box = add_shape(slide, Inches(6.8), Inches(1.7), Inches(5.8), Inches(4.8), RGBColor(0x25, 0x25, 0x40))
add_text_box(slide, Inches(7.1), Inches(1.8), Inches(5), Inches(0.5),
             "TopBrain 2025", font_size=22, color=ACCENT2, bold=True)
add_text_box(slide, Inches(7.1), Inches(2.3), Inches(5), Inches(0.4),
             "Extended Annotations — Same Scans", font_size=13, color=MID_GRAY)
add_bullet_slide_content(slide, Inches(7.1), Inches(2.8), Inches(5.2), Inches(3.5), [
    (" (same 25 CT scans as TopCoW subset)", "25 CTA volumes"),
    (" vessel classes — 3x more coverage", "40"),
    (": + distal branches, posterior fossa, small arteries, venous sinuses", "New"),
    (" ground truth covers far more vasculature", "Binarized"),
    (" for fine-tuning & transfer evaluation", "Used"),
], font_size=15)


# =============================================================================
# SLIDE 6: ARCHITECTURE
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_header(slide, 5, "Architecture: 3D U-Net")

add_image_safe(slide, os.path.join(FIG_DIR, "unet3d_architecture.png"),
               Inches(0.5), Inches(1.5), width=Inches(6.5))

specs = [
    ("Encoder", "5 stages: 32 → 64 → 128 → 256 → 512 channels"),
    ("Blocks", "2× Conv3D + InstanceNorm + LeakyReLU + residual"),
    ("Down", "Strided 2×2×2 convolutions (no max-pool)"),
    ("Up", "Transposed convolutions + skip connections"),
    ("Deep supervision", "1×1 heads at each decoder level"),
    ("Parameters", "24.9M"),
    ("Design", "nnU-Net philosophy (Isensee et al.)"),
]
for i, (key, val) in enumerate(specs):
    y = Inches(1.7 + i * 0.6)
    add_text_box(slide, Inches(7.5), y, Inches(2.5), Inches(0.5),
                 key, font_size=14, color=ACCENT, bold=True)
    add_text_box(slide, Inches(9.5), y, Inches(3.5), Inches(0.5),
                 val, font_size=13, color=WHITE)


# =============================================================================
# SLIDE 7: LOSS FUNCTIONS
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_header(slide, 6, "Loss Functions")

losses = [
    ("Dice + Cross-Entropy", "Measures voxel overlap.\nNo connectivity incentive.\nStandard nnU-Net default.",
     "L = L_CE + L_Dice", ACCENT),
    ("+ Soft clDice", "Computes Dice on skeletons of\nprediction and ground truth.\nDifferentiable soft-skeletonization\n(10 iterations of morphological ops).",
     "L = (1-α)·L_base + α·(1 - clDice)", ACCENT2),
    ("+ Skeleton Recall", "Precomputes GT skeleton on CPU.\nWeighted CE on centerline voxels.\n\"Did you find the core of\nevery vessel?\"",
     "L = (1-α)·L_base + α·(1 - recall)", GREEN_ACCENT),
]
for i, (title, desc, formula, color) in enumerate(losses):
    left = Inches(0.5 + i * 4.2)
    box = add_shape(slide, left, Inches(1.7), Inches(3.9), Inches(4.5), RGBColor(0x25, 0x25, 0x40))
    add_text_box(slide, left + Inches(0.3), Inches(1.85), Inches(3.4), Inches(0.5),
                 title, font_size=18, color=color, bold=True)
    add_shape(slide, left + Inches(0.3), Inches(2.35), Inches(2), Inches(0.03), color)
    add_text_box(slide, left + Inches(0.3), Inches(2.6), Inches(3.3), Inches(2.5),
                 desc, font_size=14, color=WHITE)
    # Formula box
    add_shape(slide, left + Inches(0.2), Inches(4.9), Inches(3.5), Inches(0.7), RGBColor(0x15, 0x15, 0x28))
    add_text_box(slide, left + Inches(0.3), Inches(4.95), Inches(3.3), Inches(0.6),
                 formula, font_size=12, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0.7), Inches(6.5), Inches(11), Inches(0.5),
             "All combined with base Dice+CE.  Topology weight α = 0.5.  Deep supervision wrapper at all decoder scales.",
             font_size=14, color=MID_GRAY)


# =============================================================================
# SLIDE 8: TRAINING SETUP
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_header(slide, 7, "Training Setup")

params = [
    ("Hardware", "8× NVIDIA A100-SXM4-40GB (AWS p4d.24xlarge)"),
    ("Training", "DDP across all 8 GPUs — identical for all 3 models"),
    ("Epochs", "300 (from-scratch)  |  150 (TopBrain fine-tuning)"),
    ("Optimizer", "AdamW (LR=1e-3, weight decay=1e-5)"),
    ("Schedule", "Cosine annealing + 10-epoch linear warmup"),
    ("Patches", "128³ voxel crops, 4 per volume, 33% foreground-centered"),
    ("Batch size", "4 per GPU  |  Mixed precision (AMP)"),
    ("Validation", "Every 25 epochs  |  Early stopping patience = 5"),
    ("Fine-tuning", "LR=1e-4, fresh optimizer, 5-epoch warmup"),
]
for i, (key, val) in enumerate(params):
    y = Inches(1.8 + i * 0.55)
    add_shape(slide, Inches(0.7), y, Inches(0.08), Inches(0.35), ACCENT)
    add_text_box(slide, Inches(1.0), y, Inches(3), Inches(0.5),
                 key, font_size=16, color=ACCENT, bold=True)
    add_text_box(slide, Inches(4.0), y, Inches(8.5), Inches(0.5),
                 val, font_size=15, color=WHITE)

# Emphasis box
box = add_shape(slide, Inches(2), Inches(6.6), Inches(9), Inches(0.6), RGBColor(0x00, 0x3D, 0x5B))
add_text_box(slide, Inches(2.3), Inches(6.62), Inches(8.5), Inches(0.55),
             "All training configuration is identical across the three runs — controlled comparison",
             font_size=15, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)


# =============================================================================
# SLIDE 9: LEARNING CURVES
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_header(slide, 8, "Training Dynamics")

add_image_safe(slide, os.path.join(FIG_DIR, "learning_curves.png"),
               Inches(1.5), Inches(1.4), height=Inches(5.5))

add_text_box(slide, Inches(8.5), Inches(2.0), Inches(4.2), Inches(4.5),
             "All three models converge to\nsimilar training loss (~0.10-0.13)\nand validation Dice (~0.84-0.86).\n\n"
             "clDice and Dice+CE curves\nare nearly identical.\n\n"
             "Skeleton Recall runs at slightly\nhigher loss — expected from\nthe additional centerline penalty.\n\n"
             "Stable convergence, no NaN\nbatches across all configurations.",
             font_size=15, color=WHITE)


# =============================================================================
# SLIDE 10: GLOBAL RESULTS
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_header(slide, 9, "Global Results — From Scratch")

add_image_safe(slide, os.path.join(FIG_DIR, "global_comparison.png"),
               Inches(0.3), Inches(1.5), width=Inches(7))

# Results callout
results = [
    ("clDice wins aggregate:", ""),
    ("  DSC 0.864  |  clDice 0.916", ""),
    ("  Lowest Betti-0 error: 3.72", ""),
    ("", ""),
    ("Baseline close behind:", ""),
    ("  DSC 0.858  |  clDice 0.907", ""),
    ("", ""),
    ("Skeleton Recall:", ""),
    ("  Lower aggregate (DSC 0.845)", ""),
    ("  High Betti-0 error: 26.84", ""),
    ("  But watch what happens next...", ""),
]
add_bullet_slide_content(slide, Inches(7.8), Inches(1.6), Inches(5), Inches(5), results,
                          font_size=15, color=WHITE)


# =============================================================================
# SLIDE 11: STRATIFIED RESULTS
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_header(slide, 10, "Stratified Results — The Vessel Gap")

add_image_safe(slide, os.path.join(FIG_DIR, "vessel_gap.png"),
               Inches(0.3), Inches(1.5), width=Inches(7))

callouts = [
    ("Skeleton Recall improves Pcom:", ""),
    ("  R-Pcom: 0.511 vs 0.459 baseline", ""),
    ("  L-Pcom: 0.479 vs 0.410 baseline", ""),
    ("  (+5 to 7 points on thinnest vessels)", ""),
    ("", ""),
    ("Smallest gap: 0.332", ""),
    ("  (vs 0.363 for both others)", ""),
    ("", ""),
    ("Hypothesis partially supported:", ""),
    ("  Modest but targeted improvement", ""),
    ("  on communicating arteries", ""),
]
add_bullet_slide_content(slide, Inches(7.8), Inches(1.6), Inches(5), Inches(5), callouts,
                          font_size=15, color=WHITE)


# =============================================================================
# SLIDE 12: TOPBRAIN — THE BIG RESULT
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_header(slide, 11, "TopBrain Fine-Tuning — The Big Result")

add_image_safe(slide, os.path.join(FIG_DIR, "topbrain_comparison.png"),
               Inches(0.3), Inches(1.5), width=Inches(7.5))

# Key finding callout
box = add_shape(slide, Inches(8), Inches(1.6), Inches(4.8), Inches(2.5), RGBColor(0x00, 0x3D, 0x5B))
add_text_box(slide, Inches(8.3), Inches(1.7), Inches(4.3), Inches(0.4),
             "SMALL ARTERIES (OA, AChA)", font_size=13, color=RED_ACCENT, bold=True)
add_text_box(slide, Inches(8.3), Inches(2.1), Inches(4.3), Inches(0.7),
             "Skeleton:  0.589", font_size=28, color=GREEN_ACCENT, bold=True)
add_text_box(slide, Inches(8.3), Inches(2.8), Inches(4.3), Inches(0.7),
             "Dice+CE:  0.000\nclDice:      0.000", font_size=18, color=RED_ACCENT, bold=True)

box2 = add_shape(slide, Inches(8), Inches(4.3), Inches(4.8), Inches(2.8), RGBColor(0x25, 0x25, 0x40))
results2 = [
    ("Posterior fossa: ", "0.626 vs 0.565"),
    ("Venous sinuses: ", "0.613 vs 0.416"),
    ("Distal branches: ", "0.752 vs 0.687"),
    ("Large CoW: ", "all ~0.82 (comparable)"),
    ("", ""),
    ("Skeleton Recall finds vessels", ""),
    ("the others cannot.", ""),
]
add_bullet_slide_content(slide, Inches(8.2), Inches(4.4), Inches(4.4), Inches(2.6), results2,
                          font_size=14, color=WHITE)


# =============================================================================
# SLIDE 13: QUALITATIVE (PLACEHOLDER)
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_header(slide, 12, "Qualitative Comparison — 3D Renderings")

# Placeholder boxes for screenshots
for i, label in enumerate(["Dice+CE + TopBrain", "Skeleton Recall + TopBrain"]):
    left = Inches(0.7 + i * 6.3)
    box = add_shape(slide, left, Inches(1.6), Inches(5.8), Inches(4.5), RGBColor(0x25, 0x25, 0x40))
    add_text_box(slide, left + Inches(0.5), Inches(3.2), Inches(4.8), Inches(1),
                 "[SurgicalAR Screenshot]", font_size=20, color=MID_GRAY,
                 alignment=PP_ALIGN.CENTER)
    color = ACCENT if i == 0 else GREEN_ACCENT
    add_text_box(slide, left, Inches(6.2), Inches(5.8), Inches(0.5),
                 label, font_size=18, color=color, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0.7), Inches(6.8), Inches(11.5), Inches(0.5),
             "Same patient  |  Same viewing angle  |  Same fine-tuning protocol  |  Different loss during pretraining",
             font_size=14, color=MID_GRAY, alignment=PP_ALIGN.CENTER)


# =============================================================================
# SLIDE 14: DISCUSSION
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_header(slide, 13, "Discussion — What We Learned")

# Two-column layout
col_data = [
    ("clDice", ACCENT, [
        "Best aggregate model",
        "DSC 0.864 | B0 error 3.72",
        "Improves topological precision",
        "Fewer disconnected fragments",
    ]),
    ("Skeleton Recall", GREEN_ACCENT, [
        "Best thin-vessel transfer",
        "Pcom +5-7 pts | Small arteries 0.589",
        "Improves topological recall",
        "Over-segments (B0 error 26.84)",
    ]),
]

for i, (title, color, bullets) in enumerate(col_data):
    left = Inches(0.7 + i * 6.3)
    box = add_shape(slide, left, Inches(1.7), Inches(5.8), Inches(3.2), RGBColor(0x25, 0x25, 0x40))
    add_text_box(slide, left + Inches(0.3), Inches(1.8), Inches(5), Inches(0.5),
                 title, font_size=20, color=color, bold=True)
    add_shape(slide, left + Inches(0.3), Inches(2.3), Inches(2), Inches(0.03), color)
    for j, b in enumerate(bullets):
        add_text_box(slide, left + Inches(0.3), Inches(2.5 + j * 0.55), Inches(5.2), Inches(0.5),
                     f"  {b}", font_size=15, color=WHITE)

# Key insight
box = add_shape(slide, Inches(0.7), Inches(5.2), Inches(11.9), Inches(1.8), RGBColor(0x00, 0x3D, 0x5B))
add_text_box(slide, Inches(1.0), Inches(5.3), Inches(11), Inches(0.4),
             "KEY INSIGHT", font_size=12, color=HIGHLIGHT, bold=True)
add_text_box(slide, Inches(1.0), Inches(5.7), Inches(11), Inches(1.0),
             "TopBrain fine-tuning drives comprehensive vessel coverage.\n"
             "The loss function during pretraining shapes how well that coverage transfers to the thinnest structures.",
             font_size=17, color=WHITE)


# =============================================================================
# SLIDE 15: LIMITATIONS & FUTURE
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_header(slide, 14, "Limitations & Future Work")

limitations = [
    (" persists (>0.33 DSC across all models)", "Large-small vessel gap"),
    (" — Betti-0 error of 26.84 needs post-processing", "Skeleton Recall over-fragments"),
    (" clips bone (800-2000 HU) to same range as vessels (200-400 HU)", "HU window [0, 600]"),
]
add_text_box(slide, Inches(0.7), Inches(1.7), Inches(5), Inches(0.5),
             "Limitations", font_size=20, color=RED_ACCENT, bold=True)
add_bullet_slide_content(slide, Inches(0.7), Inches(2.3), Inches(11), Inches(2.5), limitations,
                          font_size=16, color=WHITE, spacing=Pt(10))

future = [
    (" to separate bone from vessels — expected highest-impact fix", "Wider HU window [0, 1500]"),
    (" to focus on under-represented thin vessels", "Class-balanced sampling"),
    (" filtering to remove small disconnected false positives", "Connected component"),
    (" from 50 model-assisted labeled CTA volumes (uploaded to RedBrick AI)", "Retraining on corrected labels"),
]
add_text_box(slide, Inches(0.7), Inches(4.5), Inches(5), Inches(0.5),
             "Future Work", font_size=20, color=GREEN_ACCENT, bold=True)
add_bullet_slide_content(slide, Inches(0.7), Inches(5.1), Inches(11), Inches(2.5), future,
                          font_size=16, color=WHITE, spacing=Pt(10))


# =============================================================================
# SLIDE 16: CONCLUSION
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_header(slide, 15, "Conclusion")

conclusions = [
    ("Controlled ablation on identical hardware", "— 3 loss functions, same 3D U-Net, same 8× A100 config"),
    ("clDice = best aggregate model", "— DSC 0.864, lowest Betti-0 error, best topological precision"),
    ("Skeleton Recall = best thin-vessel transfer", "— +5-7 pts on Pcom, 0.589 on small arteries vs zero"),
    ("TopBrain fine-tuning drives coverage", "— loss function shapes transfer quality to thinnest structures"),
    ("Hypothesis partially supported", "— modest CoW improvement, dramatic effect on truly small vessels"),
]
for i, (main, sub) in enumerate(conclusions):
    y = Inches(1.8 + i * 1.0)
    add_shape(slide, Inches(0.9), y + Inches(0.05), Inches(0.08), Inches(0.35), ACCENT)
    add_text_box(slide, Inches(1.2), y, Inches(11), Inches(0.5),
                 main, font_size=19, color=WHITE, bold=True)
    add_text_box(slide, Inches(1.5), y + Inches(0.4), Inches(10.5), Inches(0.5),
                 sub, font_size=14, color=MID_GRAY)


# =============================================================================
# SLIDE 17: THANK YOU
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.12), ACCENT)
add_shape(slide, Inches(0), H - Inches(0.12), W, Inches(0.12), ACCENT)

add_text_box(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(1),
             "Thank You", font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_shape(slide, Inches(5.5), Inches(3.8), Inches(2.5), Inches(0.04), ACCENT)

add_text_box(slide, Inches(1.5), Inches(4.2), Inches(10), Inches(0.6),
             "Questions?", font_size=28, color=ACCENT, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(5.5), Inches(10), Inches(1),
             "Samih Tharwat Amer\nWhiting School of Engineering  |  Johns Hopkins University",
             font_size=16, color=MID_GRAY, alignment=PP_ALIGN.CENTER)


# =============================================================================
# SAVE
# =============================================================================
out_path = os.path.join(REPORT_DIR, "presentation.pptx")
prs.save(out_path)
print(f"Presentation saved to: {out_path}")
print(f"Slides: {len(prs.slides)}")
